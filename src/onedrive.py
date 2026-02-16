"""
OneDrive API client for MCP server.
Extracted from drive_indexer.ipynb
Uses Microsoft Graph API with device flow authentication.
"""

import io
import json
import msal
import requests
from pathlib import Path
from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# Text-based file extensions we can read directly
TEXT_EXTENSIONS = ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.py', '.js', '.ts', '.yaml', '.yml']

# Office file extensions
OFFICE_EXTENSIONS = ['.docx', '.xlsx', '.pptx']

# Azure App Registration config
CLIENT_ID = "c7b60d92-d23f-474b-9708-fb8890be59e3"
TENANT_ID = "c57d3288-4e87-42e2-bd6c-fb6f632680c3"
AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
SCOPES = ["Files.Read.All", "User.Read"]

# Token cache file
BASE_DIR = Path(__file__).parent.parent
TOKEN_CACHE_FILE = BASE_DIR / 'onedrive_token_cache.json'

# Global token storage
_access_token = None


def _get_token_cache():
    """Load or create MSAL token cache."""
    cache = msal.SerializableTokenCache()
    if TOKEN_CACHE_FILE.exists():
        cache.deserialize(TOKEN_CACHE_FILE.read_text())
    return cache


def _save_token_cache(cache):
    """Save MSAL token cache to file."""
    if cache.has_state_changed:
        TOKEN_CACHE_FILE.write_text(cache.serialize())


def authenticate(force_new: bool = False) -> str:
    """
    Authenticate with Microsoft Graph API.
    Returns access token.
    Uses device flow which requires user to visit a URL and enter a code.
    """
    global _access_token

    cache = _get_token_cache()
    app = msal.PublicClientApplication(CLIENT_ID, authority=AUTHORITY, token_cache=cache)

    # Try to get token silently from cache
    accounts = app.get_accounts()
    if accounts and not force_new:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            _save_token_cache(cache)
            _access_token = result["access_token"]
            return _access_token

    # Need to do device flow
    flow = app.initiate_device_flow(scopes=SCOPES)
    if "user_code" not in flow:
        raise Exception("Failed to initiate device flow")

    print("\n" + "=" * 50)
    print("OneDrive Authentication Required")
    print("=" * 50)
    print(f"Go to: {flow['verification_uri']}")
    print(f"Enter code: {flow['user_code']}")
    print("=" * 50 + "\n")

    result = app.acquire_token_by_device_flow(flow)

    if "access_token" not in result:
        raise Exception(f"Authentication failed: {result.get('error_description', 'Unknown error')}")

    _save_token_cache(cache)
    _access_token = result["access_token"]
    return _access_token


def _get_headers():
    """Get authorization headers for API calls."""
    global _access_token
    if not _access_token:
        authenticate()
    return {"Authorization": f"Bearer {_access_token}"}


def _api_get(url: str) -> dict:
    """Make GET request to Graph API."""
    response = requests.get(url, headers=_get_headers())
    if response.status_code == 401:
        # Token expired, re-auth and retry
        authenticate(force_new=True)
        response = requests.get(url, headers=_get_headers())
    response.raise_for_status()
    return response.json()


def list_root_items() -> list[dict]:
    """List items in OneDrive root folder."""
    data = _api_get("https://graph.microsoft.com/v1.0/me/drive/root/children")
    return data.get("value", [])


def list_folder_by_path(folder_path: str) -> list[dict]:
    """List items in a folder by path (e.g., 'Documents/Projects')."""
    encoded_path = folder_path.replace(" ", "%20")
    url = f"https://graph.microsoft.com/v1.0/me/drive/root:/{encoded_path}:/children"
    data = _api_get(url)
    return data.get("value", [])


def list_folder_by_id(folder_id: str) -> list[dict]:
    """List items in a folder by ID."""
    url = f"https://graph.microsoft.com/v1.0/me/drive/items/{folder_id}/children"
    data = _api_get(url)
    return data.get("value", [])


def list_files_recursive(folder_id: str, path: str = "") -> list[dict]:
    """Recursively list all files in a folder and subfolders."""
    all_files = []
    items = list_folder_by_id(folder_id)

    for item in items:
        current_path = f"{path}/{item['name']}" if path else item['name']
        if item.get("folder"):
            # Recurse into subfolder
            all_files.extend(list_files_recursive(item["id"], current_path))
        else:
            all_files.append({
                'id': item['id'],
                'name': item['name'],
                'path': current_path,
                'webUrl': item.get('webUrl', ''),
                'lastModifiedDateTime': item.get('lastModifiedDateTime', ''),
                'size': item.get('size', 0)
            })

    return all_files


def search_files(query: str, max_results: int = 20) -> list[dict]:
    """Search for files by name."""
    url = f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{query}')?$top={max_results}"
    data = _api_get(url)
    return data.get("value", [])


def get_file_by_path(file_path: str) -> dict:
    """Get file metadata by path."""
    encoded_path = file_path.replace(" ", "%20")
    url = f"https://graph.microsoft.com/v1.0/me/drive/root:/{encoded_path}"
    return _api_get(url)


def _download_file_content(download_url: str) -> bytes:
    """Download file content from a download URL."""
    response = requests.get(download_url, headers=_get_headers())
    if response.status_code == 401:
        authenticate(force_new=True)
        response = requests.get(download_url, headers=_get_headers())
    response.raise_for_status()
    return response.content


def _extract_docx_text(content: bytes) -> str:
    """Extract text from a .docx file."""
    buffer = io.BytesIO(content)
    doc = DocxDocument(buffer)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return '\n\n'.join(paragraphs)


def _extract_xlsx_text(content: bytes) -> str:
    """Extract text from a .xlsx file as CSV-like format."""
    buffer = io.BytesIO(content)

    # Try different modes if one fails
    try:
        wb = load_workbook(buffer, read_only=True, data_only=True)
    except Exception:
        buffer.seek(0)
        try:
            wb = load_workbook(buffer, data_only=True)
        except Exception:
            buffer.seek(0)
            wb = load_workbook(buffer)

    all_sheets = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            # Convert None to empty string and join with comma
            row_str = ','.join(str(cell) if cell is not None else '' for cell in row)
            if row_str.strip(','):  # Skip completely empty rows
                rows.append(row_str)

        if rows:
            sheet_text = f"=== Sheet: {sheet_name} ===\n" + '\n'.join(rows)
            all_sheets.append(sheet_text)

    wb.close()
    return '\n\n'.join(all_sheets) if all_sheets else "Spreadsheet is empty or could not be parsed."


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from a .pptx file."""
    buffer = io.BytesIO(content)
    prs = Presentation(buffer)

    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"=== Slide {i} ==="]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text)
        if len(slide_parts) > 1:  # Has content beyond just the header
            slides_text.append('\n'.join(slide_parts))

    return '\n\n'.join(slides_text)


def read_file_by_path(file_path: str, max_chars: int = 100000) -> str:
    """
    Read file content by path.

    Supports:
    - Text files (.txt, .md, .csv, .json, etc.)
    - PDFs (.pdf)
    - Word documents (.docx)
    - Excel spreadsheets (.xlsx)
    - PowerPoint presentations (.pptx)
    """
    try:
        # Get file metadata including download URL
        file_info = get_file_by_path(file_path)

        if 'folder' in file_info:
            return f"Cannot read folder: {file_path}"

        file_name = file_info.get('name', file_path)
        download_url = file_info.get('@microsoft.graph.downloadUrl')

        if not download_url:
            return f"No download URL available for: {file_path}"

        # Download the file content
        content = _download_file_content(download_url)

        # Determine file type by extension
        ext = Path(file_name).suffix.lower()

        # PDF → extract text
        if ext == '.pdf':
            buffer = io.BytesIO(content)
            reader = PdfReader(buffer)
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            text = '\n'.join(text_parts)

        # Word document → extract text
        elif ext == '.docx':
            text = _extract_docx_text(content)

        # Excel spreadsheet → extract as CSV-like
        elif ext == '.xlsx':
            text = _extract_xlsx_text(content)

        # PowerPoint → extract text
        elif ext == '.pptx':
            text = _extract_pptx_text(content)

        # Text-based files → decode as UTF-8
        elif ext in TEXT_EXTENSIONS or file_info.get('file', {}).get('mimeType', '').startswith('text/'):
            text = content.decode('utf-8')

        else:
            return f"Cannot read file type: {ext} ({file_name})"

        # Truncate if too long
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n... [truncated, {len(text)} total chars]"

        return text

    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 404:
            return f"File not found: {file_path}"
        return f"Error reading file: {str(e)}"
    except Exception as e:
        return f"Error reading file: {str(e)}"


def read_file_by_id(file_id: str, max_chars: int = 100000) -> str:
    """
    Read file content by ID.

    Supports:
    - Text files (.txt, .md, .csv, .json, etc.)
    - PDFs (.pdf)
    - Word documents (.docx)
    - Excel spreadsheets (.xlsx)
    - PowerPoint presentations (.pptx)
    """
    try:
        # Get file metadata
        url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}"
        file_info = _api_get(url)

        if 'folder' in file_info:
            return f"Cannot read folder: {file_info.get('name', file_id)}"

        file_name = file_info.get('name', '')
        download_url = file_info.get('@microsoft.graph.downloadUrl')

        if not download_url:
            return f"No download URL available for: {file_name}"

        # Download the file content
        content = _download_file_content(download_url)

        # Determine file type by extension
        ext = Path(file_name).suffix.lower()

        # PDF → extract text
        if ext == '.pdf':
            buffer = io.BytesIO(content)
            reader = PdfReader(buffer)
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            text = '\n'.join(text_parts)

        # Word document → extract text
        elif ext == '.docx':
            text = _extract_docx_text(content)

        # Excel spreadsheet → extract as CSV-like
        elif ext == '.xlsx':
            text = _extract_xlsx_text(content)

        # PowerPoint → extract text
        elif ext == '.pptx':
            text = _extract_pptx_text(content)

        # Text-based files → decode as UTF-8
        elif ext in TEXT_EXTENSIONS or file_info.get('file', {}).get('mimeType', '').startswith('text/'):
            text = content.decode('utf-8')

        else:
            return f"Cannot read file type: {ext} ({file_name})"

        # Truncate if too long
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n... [truncated, {len(text)} total chars]"

        return text

    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 404:
            return f"File not found: {file_id}"
        return f"Error reading file: {str(e)}"
    except Exception as e:
        return f"Error reading file: {str(e)}"


# Quick test when run directly
if __name__ == "__main__":
    print("Testing OneDrive connection...")
    token = authenticate()
    print("✓ Authenticated successfully")

    items = list_root_items()
    print(f"✓ Found {len(items)} items in root:")
    for item in items:
        icon = "📁" if item.get("folder") else "📄"
        print(f"  {icon} {item['name']}")
